import os
import pandas as pd
import PyPDF2
import pdfplumber
from pptx import Presentation
import chardet
import re
from datetime import datetime
from dateutil.parser import parse
from anthropic import Anthropic
import json

class MediaFileProcessor:
    def __init__(self, anthropic_api_key):
        self.anthropic = Anthropic(api_key=anthropic_api_key)
        
    def process_file(self, file_path, filename):
        """Process a file based on its extension and return standardized data."""
        file_ext = os.path.splitext(filename)[1].lower()
        
        try:
            if file_ext in ['.xlsx', '.xls']:
                return self.process_excel(file_path)
            elif file_ext == '.pdf':
                return self.process_pdf(file_path)
            elif file_ext == '.pptx':
                return self.process_pptx(file_path)
            else:
                raise ValueError(f"Unsupported file type: {file_ext}")
        except Exception as e:
            print(f"Error processing file {filename}: {str(e)}")
            return []
    
    def process_excel(self, file_path):
        """Process Excel files with intelligent column mapping."""
        try:
            # Read all sheets
            excel_file = pd.ExcelFile(file_path)
            all_data = []
            
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(file_path, sheet_name=sheet_name)
                
                # Skip empty sheets
                if df.empty:
                    continue
                
                # Detect and map columns using AI
                mapped_data = self.intelligent_column_mapping(df)
                all_data.extend(mapped_data)
            
            return all_data
        except Exception as e:
            print(f"Error processing Excel file: {str(e)}")
            return []
    
    def intelligent_column_mapping(self, df):
        """Use AI to intelligently map Excel columns to standard format."""
        # Get column headers and sample data
        headers = list(df.columns)
        sample_rows = df.head(3).to_dict('records')
        
        # Create prompt for Claude to analyze the structure
        prompt = f"""
        Analyze this Excel data structure and map the columns to standard media coverage fields.
        
        Column headers: {headers}
        Sample data: {sample_rows}
        
        Please identify which columns correspond to:
        - headline/title (article headline or title)
        - date (publication date)
        - outlet/source (media outlet or publication name)
        - description/content (article content or summary)
        - url (article URL if available)
        - author (author name if available)
        - topic/category (topic or category if available)
        
        Return a JSON object with the mapping:
        {{
            "headline": "column_name_or_null",
            "date": "column_name_or_null", 
            "outlet": "column_name_or_null",
            "description": "column_name_or_null",
            "url": "column_name_or_null",
            "author": "column_name_or_null",
            "topic": "column_name_or_null"
        }}
        
        If a field cannot be mapped to any column, use null.
        """
        
        try:
            response = self.anthropic.messages.create(
                model="claude-3-haiku-20240307",
                max_tokens=500,
                messages=[{"role": "user", "content": prompt}]
            )
            
            # Extract JSON from response
            response_text = response.content[0].text
            json_match = re.search(r'\{.*\}', response_text, re.DOTALL)
            
            if json_match:
                mapping = json.loads(json_match.group(0))
            else:
                # Fallback to basic mapping
                mapping = self.basic_column_mapping(headers)
                
        except Exception as e:
            print(f"AI mapping failed, using basic mapping: {str(e)}")
            mapping = self.basic_column_mapping(headers)
        
        # Convert DataFrame to standardized format
        return self.convert_to_standard_format(df, mapping)
    
    def basic_column_mapping(self, headers):
        """Fallback basic column mapping based on common patterns."""
        mapping = {
            "headline": None,
            "date": None,
            "outlet": None,
            "description": None,
            "url": None,
            "author": None,
            "topic": None
        }
        
        headers_lower = [h.lower() for h in headers]
        
        # Common patterns for each field
        patterns = {
            "headline": ["headline", "title", "article title", "story", "subject"],
            "date": ["date", "published", "publish date", "article date", "timestamp"],
            "outlet": ["outlet", "publication", "source", "media outlet", "publisher"],
            "description": ["description", "content", "summary", "text", "body"],
            "url": ["url", "link", "article url", "web link"],
            "author": ["author", "writer", "journalist", "reporter"],
            "topic": ["topic", "category", "subject", "theme", "tag"]
        }
        
        for field, keywords in patterns.items():
            for keyword in keywords:
                for i, header in enumerate(headers_lower):
                    if keyword in header:
                        mapping[field] = headers[i]
                        break
                if mapping[field]:
                    break
        
        return mapping
    
    def convert_to_standard_format(self, df, mapping):
        """Convert DataFrame to standardized article format."""
        articles = []
        
        for _, row in df.iterrows():
            try:
                article = {
                    "source": {"name": self.safe_get_value(row, mapping["outlet"], "Unknown Source")},
                    "title": self.safe_get_value(row, mapping["headline"], "No Title"),
                    "description": self.safe_get_value(row, mapping["description"], ""),
                    "url": self.safe_get_value(row, mapping["url"], ""),
                    "author": self.safe_get_value(row, mapping["author"], ""),
                    "publishedAt": self.parse_date(self.safe_get_value(row, mapping["date"], "")),
                    "content": self.safe_get_value(row, mapping["description"], ""),
                    "api_source": "Local File"
                }
                
                # Skip rows without essential data
                if article["title"] and article["title"] != "No Title":
                    articles.append(article)
                    
            except Exception as e:
                print(f"Error processing row: {str(e)}")
                continue
        
        return articles
    
    def safe_get_value(self, row, column_name, default=""):
        """Safely get value from row, handling missing columns."""
        if column_name and column_name in row:
            value = row[column_name]
            if pd.isna(value):
                return default
            return str(value)
        return default
    
    def parse_date(self, date_str):
        """Parse various date formats to ISO format."""
        if not date_str or date_str == "":
            return datetime.now().strftime("%Y-%m-%dT%H:%M:%SZ")
        
        try:
            # Try to parse the date
            parsed_date = parse(str(date_str))
            return parsed_date.strftime("%Y-%m-%dT%H:%M:%SZ")
        except:
            # Return current date if parsing fails
            return datetime.now().strftime("%Y-%m-%dT%H:%M:%SZ")
    
    def process_pdf(self, file_path):
        """Extract text content from PDF files."""
        articles = []
        
        try:
            # Try pdfplumber first (better for complex layouts)
            with pdfplumber.open(file_path) as pdf:
                full_text = ""
                for page in pdf.pages:
                    text = page.extract_text()
                    if text:
                        full_text += text + "\n"
            
            # If pdfplumber fails, try PyPDF2
            if not full_text.strip():
                with open(file_path, 'rb') as file:
                    pdf_reader = PyPDF2.PdfReader(file)
                    for page in pdf_reader.pages:
                        full_text += page.extract_text() + "\n"
            
            # Use AI to structure the extracted text
            if full_text.strip():
                structured_articles = self.structure_pdf_content(full_text)
                articles.extend(structured_articles)
                
        except Exception as e:
            print(f"Error processing PDF: {str(e)}")
        
        return articles
    
    def structure_pdf_content(self, text):
        """Use AI to identify and structure articles from PDF text."""
        prompt = f"""
        Analyze this text extracted from a PDF and identify individual media coverage items or articles.
        
        Text content:
        {text[:4000]}  # Limit text to avoid token limits
        
        Please identify and extract:
        1. Individual articles or media mentions
        2. Headlines or titles
        3. Publication dates (if mentioned)
        4. Source publications (if mentioned)
        5. Main content or summaries
        
        Return a JSON array of articles in this format:
        [
            {{
                "title": "Article headline or title",
                "source": "Publication name or Unknown",
                "date": "Date in YYYY-MM-DD format or null",
                "content": "Main article content or summary",
                "description": "Brief description or summary"
            }}
        ]
        
        If the text appears to be one continuous document, create a single article entry.
        If you can identify multiple distinct articles, create separate entries.
        """
        
        try:
            response = self.anthropic.messages.create(
                model="claude-3-haiku-20240307",
                max_tokens=1000,
                messages=[{"role": "user", "content": prompt}]
            )
            
            response_text = response.content[0].text
            json_match = re.search(r'\[.*\]', response_text, re.DOTALL)
            
            if json_match:
                articles_data = json.loads(json_match.group(0))
                return self.convert_pdf_to_standard_format(articles_data)
            else:
                # Fallback: create single article from all text
                return self.create_single_article_from_text(text)
                
        except Exception as e:
            print(f"AI structuring failed: {str(e)}")
            return self.create_single_article_from_text(text)
    
    def convert_pdf_to_standard_format(self, articles_data):
        """Convert AI-structured PDF data to standard format."""
        articles = []
        
        for article_data in articles_data:
            try:
                article = {
                    "source": {"name": article_data.get("source", "PDF Document")},
                    "title": article_data.get("title", "PDF Content"),
                    "description": article_data.get("description", article_data.get("content", "")[:200]),
                    "url": "",
                    "author": "",
                    "publishedAt": self.parse_date(article_data.get("date", "")),
                    "content": article_data.get("content", ""),
                    "api_source": "Local File"
                }
                articles.append(article)
            except Exception as e:
                print(f"Error converting PDF article: {str(e)}")
                continue
        
        return articles
    
    def create_single_article_from_text(self, text):
        """Create a single article from unstructured text."""
        return [{
            "source": {"name": "PDF Document"},
            "title": "PDF Content Analysis",
            "description": text[:200] + "..." if len(text) > 200 else text,
            "url": "",
            "author": "",
            "publishedAt": datetime.now().strftime("%Y-%m-%dT%H:%M:%SZ"),
            "content": text,
            "api_source": "Local File"
        }]
    
    def process_pptx(self, file_path):
        """Extract text content from PowerPoint files."""
        articles = []
        
        try:
            prs = Presentation(file_path)
            slides_content = []
            
            for i, slide in enumerate(prs.slides):
                slide_text = ""
                slide_title = f"Slide {i + 1}"
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        if not slide_title or slide_title == f"Slide {i + 1}":
                            # Use first text as title if no title found
                            slide_title = shape.text.strip()[:100]
                        slide_text += shape.text + "\n"
                
                if slide_text.strip():
                    slides_content.append({
                        "title": slide_title,
                        "content": slide_text.strip()
                    })
            
            # Convert slides to articles
            for slide_data in slides_content:
                article = {
                    "source": {"name": "PowerPoint Presentation"},
                    "title": slide_data["title"],
                    "description": slide_data["content"][:200] + "..." if len(slide_data["content"]) > 200 else slide_data["content"],
                    "url": "",
                    "author": "",
                    "publishedAt": datetime.now().strftime("%Y-%m-%dT%H:%M:%SZ"),
                    "content": slide_data["content"],
                    "api_source": "Local File"
                }
                articles.append(article)
                
        except Exception as e:
            print(f"Error processing PowerPoint: {str(e)}")
        
        return articles

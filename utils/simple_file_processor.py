import os
import re
from datetime import datetime
from anthropic import Anthropic

class SimpleMediaFileProcessor:
    def __init__(self, anthropic_api_key):
        self.anthropic = Anthropic(api_key=anthropic_api_key)
    
    def process_file(self, file_path, original_filename):
        """Process uploaded files and extract media coverage data."""
        try:
            file_extension = original_filename.lower().split('.')[-1]
            
            if file_extension in ['xlsx', 'xls']:
                return self._process_excel_simple(file_path, original_filename)
            elif file_extension == 'pdf':
                return self._process_pdf_simple(file_path, original_filename)
            elif file_extension == 'pptx':
                return self._process_pptx_simple(file_path, original_filename)
            else:
                print(f"Unsupported file type: {file_extension}")
                return []
                
        except Exception as e:
            print(f"Error processing file {original_filename}: {str(e)}")
            return []
    
    def _process_excel_simple(self, file_path, filename):
        """Simple Excel processing without pandas."""
        try:
            import openpyxl
            
            workbook = openpyxl.load_workbook(file_path)
            articles = []
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                
                # Try to find header row
                headers = []
                data_start_row = 1
                
                for row in range(1, min(6, sheet.max_row + 1)):  # Check first 5 rows for headers
                    row_values = [cell.value for cell in sheet[row] if cell.value]
                    if any(header in str(cell.value).lower() if cell.value else '' 
                          for cell in sheet[row] 
                          for header in ['title', 'headline', 'article', 'date', 'source', 'publication']):
                        headers = [cell.value for cell in sheet[row]]
                        data_start_row = row + 1
                        break
                
                if not headers:
                    # No clear headers found, assume first row is data
                    headers = [f"Column_{i+1}" for i in range(sheet.max_column)]
                    data_start_row = 1
                
                # Process data rows
                for row in range(data_start_row, sheet.max_row + 1):
                    row_data = {}
                    has_content = False
                    
                    for col, header in enumerate(headers, 1):
                        if col <= sheet.max_column:
                            cell_value = sheet.cell(row=row, column=col).value
                            if cell_value:
                                row_data[str(header)] = str(cell_value)
                                has_content = True
                    
                    if has_content:
                        article = self._extract_article_from_row(row_data, filename)
                        if article:
                            articles.append(article)
            
            return articles[:100]  # Limit to 100 articles
            
        except Exception as e:
            print(f"Error processing Excel file: {str(e)}")
            return []
    
    def _process_pdf_simple(self, file_path, filename):
        """Simple PDF processing."""
        try:
            import pdfplumber
            
            articles = []
            text_content = ""
            
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages[:10]:  # Limit to first 10 pages
                    page_text = page.extract_text()
                    if page_text:
                        text_content += page_text + "\n"
            
            if text_content:
                # Try to extract articles using AI
                articles = self._extract_articles_from_text(text_content, filename)
            
            return articles
            
        except Exception as e:
            print(f"Error processing PDF file: {str(e)}")
            return []
    
    def _process_pptx_simple(self, file_path, filename):
        """Simple PowerPoint processing."""
        try:
            from pptx import Presentation
            
            articles = []
            text_content = ""
            
            prs = Presentation(file_path)
            
            for slide in prs.slides:
                slide_text = ""
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text += shape.text + "\n"
                
                if slide_text.strip():
                    text_content += slide_text + "\n\n"
            
            if text_content:
                articles = self._extract_articles_from_text(text_content, filename)
            
            return articles
            
        except Exception as e:
            print(f"Error processing PowerPoint file: {str(e)}")
            return []
    
    def _extract_article_from_row(self, row_data, filename):
        """Extract article information from a spreadsheet row."""
        try:
            # Try to identify title/headline
            title = ""
            for key in row_data.keys():
                if any(word in key.lower() for word in ['title', 'headline', 'article', 'subject']):
                    title = row_data[key]
                    break
            
            if not title:
                # Use first non-empty value as title
                title = next(iter(row_data.values()), "")
            
            # Try to identify date
            date_str = ""
            for key in row_data.keys():
                if any(word in key.lower() for word in ['date', 'time', 'published']):
                    date_str = row_data[key]
                    break
            
            # Try to identify source
            source = ""
            for key in row_data.keys():
                if any(word in key.lower() for word in ['source', 'publication', 'outlet', 'media']):
                    source = row_data[key]
                    break
            
            if not source:
                source = filename
            
            # Create article object
            if title and len(title.strip()) > 5:
                return {
                    'title': title[:200],  # Limit title length
                    'description': ' '.join([v for k, v in row_data.items() if k != title])[:300],
                    'publishedAt': self._parse_date(date_str),
                    'source': {'name': source[:50]},
                    'url': f"file://{filename}",
                    'sentiment': 0  # Will be calculated later
                }
            
            return None
            
        except Exception as e:
            print(f"Error extracting article from row: {str(e)}")
            return None
    
    def _extract_articles_from_text(self, text_content, filename):
        """Use AI to extract articles from text content."""
        try:
            # Limit text length for API call
            if len(text_content) > 10000:
                text_content = text_content[:10000]
            
            prompt = f"""
            Extract media coverage information from the following text. 
            Return a list of articles with title, description, and any dates mentioned.
            Format as simple text with each article separated by "---".
            
            Text:
            {text_content}
            """
            
            response = self.anthropic.messages.create(
                model="claude-3-haiku-20240307",
                max_tokens=1000,
                messages=[{
                    "role": "user",
                    "content": prompt
                }]
            )
            
            # Parse AI response into articles
            articles = []
            response_text = response.content[0].text
            
            # Simple parsing - split by "---" and extract basic info
            article_blocks = response_text.split("---")
            
            for i, block in enumerate(article_blocks[:10]):  # Limit to 10 articles
                block = block.strip()
                if len(block) > 20:  # Minimum content length
                    lines = block.split('\n')
                    title = lines[0][:200] if lines else f"Article {i+1}"
                    description = ' '.join(lines[1:])[:300] if len(lines) > 1 else block[:300]
                    
                    articles.append({
                        'title': title,
                        'description': description,
                        'publishedAt': datetime.now().isoformat(),
                        'source': {'name': filename},
                        'url': f"file://{filename}",
                        'sentiment': 0
                    })
            
            return articles
            
        except Exception as e:
            print(f"Error extracting articles from text: {str(e)}")
            # Fallback: create a single article from the text
            return [{
                'title': f"Content from {filename}",
                'description': text_content[:300],
                'publishedAt': datetime.now().isoformat(),
                'source': {'name': filename},
                'url': f"file://{filename}",
                'sentiment': 0
            }]
    
    def _parse_date(self, date_str):
        """Parse date string into ISO format."""
        if not date_str:
            return datetime.now().isoformat()
        
        try:
            # Try common date formats
            for fmt in ['%Y-%m-%d', '%m/%d/%Y', '%d/%m/%Y', '%Y-%m-%d %H:%M:%S']:
                try:
                    return datetime.strptime(str(date_str), fmt).isoformat()
                except:
                    continue
            
            # If parsing fails, return current date
            return datetime.now().isoformat()
            
        except:
            return datetime.now().isoformat()
